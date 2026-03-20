"""
Accessibility Engine
Refactored from pptx_accessibility_fixer.py to support interactive web workflow.
Separates scanning from fixing, supporting per-issue approval.
"""

import os, re, uuid, base64, subprocess, shutil
from pathlib import Path
from dataclasses import dataclass, asdict
from typing import List, Optional, Dict
from lxml import etree

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER
from pptx.oxml.ns import qn

# ── NAMESPACES ─────────────────────────────────────────────────────────────────
PML = "http://schemas.openxmlformats.org/presentationml/2006/main"
DML = "http://schemas.openxmlformats.org/drawingml/2006/main"

def atag(name): return f"{{{DML}}}{name}"

DEFAULT_LANGUAGE = "en-US"
FINE_PRINT_PT    = 18
DARK_FALLBACK    = "1A1A1A"
LIGHT_THRESHOLD  = 0.55

# ── CHECK REGISTRY ─────────────────────────────────────────────────────────────
CHECKS = {
    # always_auto: applied silently, no toggle
    # always_human: flagged only, never auto-applied
    # default_auto: default toggle state for configurable checks
    "presentation_title":   {"label": "Presentation metadata title",       "category": "Structure", "always_auto": True},
    "language_tags":        {"label": "Language tags on text runs",        "category": "Structure", "always_auto": True},
    "trailing_empty_lines": {"label": "Trailing empty paragraphs",         "category": "Structure", "always_auto": True},
    "empty_textboxes":      {"label": "Empty text boxes",                  "category": "Structure", "always_auto": True},
    "table_unmerge":        {"label": "Unmerge merged table cells",        "category": "Tables",    "always_auto": True},
    "table_empty_cells":    {"label": "Fill empty table cells with —",     "category": "Tables",    "always_auto": True},

    "slide_titles":         {"label": "Missing / empty slide titles",      "category": "Structure", "default_auto": True},
    "duplicate_titles":     {"label": "Duplicate slide titles",            "category": "Structure", "default_auto": True},
    "broken_lists":         {"label": "Broken list formatting",            "category": "Content",   "default_auto": True},
    "image_alt_text":       {"label": "Image alt text (AI-generated)",     "category": "Images",    "default_auto": True},
    "table_descriptions":   {"label": "Table descriptions (alt text)",     "category": "Tables",    "default_auto": True},
    "color_contrast":       {"label": "Explicit low-contrast text colors", "category": "Visual",    "default_auto": True},

    "fine_print":           {"label": "Fine print text (< 18pt)",          "category": "Visual",    "always_human": True},
    "theme_contrast":       {"label": "Theme-inherited color contrast",    "category": "Visual",    "always_human": True},
}


# ── HELPERS ────────────────────────────────────────────────────────────────────

def sanitize(s):
    if not s: return ""
    return re.sub(r'[\x00-\x08\x0B\x0C\x0E-\x1F\x7F]', '', str(s)).strip()

def xml_escape(s):
    return sanitize(s).replace("&","&amp;").replace("<","&lt;").replace(">","&gt;").replace('"','&quot;')

def raw_text(el):
    return sanitize("".join((t.text or "") for t in el.iter(atag("t"))))

def _find_cNvPr(shape):
    el = shape._element
    for parent_tag in ["p:nvPicPr","p:nvSpPr","p:nvGraphicFramePr","p:nvGrpSpPr","p:nvCxnSpPr"]:
        parent = el.find(qn(parent_tag))
        if parent is not None:
            node = parent.find(qn("p:cNvPr"))
            if node is not None: return node
    for d in el.iter():
        local = d.tag.split("}")[-1] if "}" in d.tag else d.tag
        if local == "cNvPr": return d
    return None

def get_alt_text(shape):
    node = _find_cNvPr(shape)
    if node is not None:
        return (node.get("title","") or node.get("descr","")).strip()
    return ""

def set_alt_text(shape, title_text, descr_text=""):
    node = _find_cNvPr(shape)
    if node is not None:
        node.set("title", sanitize(title_text))
        node.set("descr", sanitize(descr_text or title_text))
        return True
    return False

def _lum(r, g, b):
    def ch(c):
        c /= 255.0
        return c / 12.92 if c <= 0.03928 else ((c + 0.055) / 1.055) ** 2.4
    return 0.2126*ch(r) + 0.7152*ch(g) + 0.0722*ch(b)

def is_too_light(hex_color):
    try:
        h = hex_color.lstrip("#")
        return _lum(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16)) > LIGHT_THRESHOLD
    except: return False

def get_run_color(run):
    try:
        rPr = run._r.find(qn("a:rPr"))
        if rPr is None: return None
        sf = rPr.find(qn("a:solidFill"))
        if sf is None: return None
        srgb = sf.find(qn("a:srgbClr"))
        return srgb.get("val") if srgb is not None else None
    except: return None

def set_run_color(run, hex_color):
    try:
        rPr = run._r.find(qn("a:rPr"))
        if rPr is None:
            rPr = etree.SubElement(run._r, qn("a:rPr"))
            run._r.insert(0, rPr)
        for sf in rPr.findall(qn("a:solidFill")): rPr.remove(sf)
        sf   = etree.SubElement(rPr, qn("a:solidFill"))
        srgb = etree.SubElement(sf,  qn("a:srgbClr"))
        srgb.set("val", hex_color.lstrip("#").upper())
        return True
    except: return False

def set_title_text(title_shape, text):
    tf = title_shape.text_frame
    if tf.paragraphs and tf.paragraphs[0].runs:
        tf.paragraphs[0].runs[0].text = text
        for r in tf.paragraphs[0].runs[1:]: r.text = ""
    else:
        p   = tf.paragraphs[0]._p
        r   = etree.SubElement(p, qn("a:r"))
        rPr = etree.SubElement(r, qn("a:rPr"))
        rPr.set("lang", DEFAULT_LANGUAGE)
        t   = etree.SubElement(r, qn("a:t"))
        t.text = text
    for extra in tf.paragraphs[1:]:
        if not raw_text(extra._p):
            try: tf._txBody.remove(extra._p)
            except: pass

def describe_table(shape, slide_title):
    try:
        table = shape.table
        rows = len(table.rows); cols = len(table.columns)
        hdrs = [c.text_frame.text.strip() for c in table.rows[0].cells] if rows > 0 else []
        hdrs = [h for h in hdrs if h]
        suffix = f": {', '.join(hdrs[:4])}" if hdrs else ""
        desc = f"Table with {rows} rows and {cols} columns{suffix}"
        return f"{slide_title} — {desc}" if slide_title else desc
    except: return f"Data table{' — ' + slide_title if slide_title else ''}"

def unmerge_table_cells(shape):
    tbl = shape._element.find(f".//{{{DML}}}tbl")
    if tbl is None: return 0
    count = 0
    for tr in tbl.findall(atag("tr")):
        for tc in tr.findall(atag("tc")):
            changed = False
            if int(tc.get("gridSpan",1)) > 1: del tc.attrib["gridSpan"]; changed = True
            if int(tc.get("rowSpan",1))  > 1: del tc.attrib["rowSpan"];  changed = True
            for attr in ("hMerge","vMerge"):
                if tc.get(attr) == "1":
                    del tc.attrib[attr]
                    txBody = tc.find(atag("txBody"))
                    if txBody is None: txBody = etree.SubElement(tc, atag("txBody"))
                    for child in list(txBody): txBody.remove(child)
                    etree.SubElement(txBody, atag("bodyPr"))
                    etree.SubElement(txBody, atag("lstStyle"))
                    p = etree.SubElement(txBody, atag("p"))
                    r = etree.SubElement(p, atag("r"))
                    t = etree.SubElement(r, atag("t")); t.text = "—"
                    changed = True
            if changed: count += 1
    return count

def fix_empty_table_cells(table):
    count = 0
    for row in table.rows:
        for cell in row.cells:
            try:
                if not cell.text_frame.text.strip():
                    p = cell.text_frame.paragraphs[0]
                    if p.runs: p.runs[0].text = "—"
                    else:
                        r = etree.SubElement(p._p, qn("a:r"))
                        t = etree.SubElement(r,    qn("a:t")); t.text = "—"
                    count += 1
            except: pass
    return count


# ── ISSUE DATACLASS ────────────────────────────────────────────────────────────

@dataclass
class Issue:
    id: str
    slide_index: int          # -1 = presentation-level
    check_type: str
    severity: str             # error | warning | info
    title: str
    description: str
    suggested_value: str = ""
    shape_name: str = ""
    status: str = "pending"   # pending | accepted | edited | skipped | auto_applied

    def to_dict(self):
        d = asdict(self)
        cfg = CHECKS.get(self.check_type, {})
        d["check_label"]   = cfg.get("label", self.check_type)
        d["always_human"]  = cfg.get("always_human", False)
        d["always_auto"]   = cfg.get("always_auto", False)
        return d


# ── SESSION CLASS ──────────────────────────────────────────────────────────────

class AccessibilitySession:

    def __init__(self, session_id: str, pptx_path: str, settings: Dict, api_key: str = None):
        self.session_id   = session_id
        self.original_path = Path(pptx_path)
        self.work_dir     = self.original_path.parent
        self.pptx_path    = self.work_dir / f"working_{self.original_path.name}"
        self.thumb_dir    = self.work_dir / "thumbnails"
        self.thumb_dir.mkdir(exist_ok=True)

        shutil.copy2(str(self.original_path), str(self.pptx_path))

        self.settings         = settings
        self.prs              = Presentation(str(self.pptx_path))
        self.issues: List[Issue] = []
        self.thumbnail_paths: List[str] = []

        self.client = None
        key = api_key or os.environ.get("ANTHROPIC_API_KEY")
        if key:
            try:
                import anthropic
                self.client = anthropic.Anthropic(api_key=key)
            except ImportError:
                pass

    # ── SETTINGS ──────────────────────────────────────────────────────────────

    def should_auto(self, check_id: str) -> bool:
        cfg = CHECKS.get(check_id, {})
        if cfg.get("always_auto"):    return True
        if cfg.get("always_human"):   return False
        return self.settings.get(check_id, cfg.get("default_auto", True))

    # ── SCAN ENTRY POINT ──────────────────────────────────────────────────────

    def scan_and_auto_fix(self):
        self._check_presentation_title()
        self._check_language_tags()
        self._check_trailing_empty_lines()
        self._check_empty_textboxes()
        self._check_slide_titles()
        self._check_broken_lists()
        self._check_image_alt_text()
        self._check_tables()
        self._check_color_contrast()
        self._check_fine_print()
        self.save_pptx()

    # ── INDIVIDUAL CHECKS ─────────────────────────────────────────────────────

    def _check_presentation_title(self):
        cur = (self.prs.core_properties.title or "").strip()
        if not cur:
            new_title = self.original_path.stem.replace("_"," ").replace("-"," ").title()
            iss = self._new("presentation_title", -1, "error",
                "Missing presentation metadata title",
                "No document title set — required by screen readers and WCAG 2.4.2.",
                suggested_value=new_title)
            if self.should_auto("presentation_title"):
                self.prs.core_properties.title = new_title
                iss.status = "auto_applied"
            self.issues.append(iss)

    def _check_language_tags(self):
        count = 0
        for slide in self.prs.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame: continue
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        rPr = run._r.find(qn("a:rPr"))
                        if rPr is None:
                            rPr = etree.SubElement(run._r, qn("a:rPr"))
                            run._r.insert(0, rPr)
                        if not rPr.get("lang"):
                            count += 1
                            if self.should_auto("language_tags"):
                                rPr.set("lang", DEFAULT_LANGUAGE)
        if count:
            iss = self._new("language_tags", -1, "error",
                f"Language tags missing on {count} text run(s)",
                "Screen readers need lang attributes to pronounce text correctly (WCAG 3.1.1).",
                suggested_value=DEFAULT_LANGUAGE)
            iss.status = "auto_applied" if self.should_auto("language_tags") else "pending"
            self.issues.append(iss)

    def _check_trailing_empty_lines(self):
        total = 0
        for slide in self.prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    total += self._remove_trailing_paras(shape.text_frame)
        if total:
            iss = self._new("trailing_empty_lines", -1, "warning",
                f"Removed {total} trailing empty paragraph(s)",
                "Empty trailing paragraphs disrupt reading order for screen readers.")
            iss.status = "auto_applied"
            self.issues.append(iss)

    def _check_empty_textboxes(self):
        for i, slide in enumerate(self.prs.slides):
            to_remove = []
            for shape in slide.shapes:
                if not shape.has_text_frame: continue
                if shape.is_placeholder:
                    ph = shape.placeholder_format.type
                    if ph in (PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE,
                              PP_PLACEHOLDER.VERTICAL_TITLE): continue
                all_text = sanitize("".join(
                    (t.text or "") for t in shape.text_frame._txBody.iter(atag("t"))
                ))
                if not all_text:
                    to_remove.append(shape)
            for shape in to_remove:
                iss = self._new("empty_textboxes", i, "warning",
                    f'Slide {i+1}: Empty text box "{shape.name}"',
                    "Empty text boxes are read aloud by screen readers as blank elements.",
                    shape_name=shape.name)
                if self.should_auto("empty_textboxes"):
                    try:
                        shape._element.getparent().remove(shape._element)
                        iss.status = "auto_applied"
                    except: iss.status = "pending"
                self.issues.append(iss)

    def _check_slide_titles(self):
        seen = {}
        for i, slide in enumerate(self.prs.slides):
            ts = self._title_shape(slide)
            current = None

            if ts is None:
                body = self._body_text(slide)
                suggested = (self._ai_title(i, body) or self._fallback_title(body) or f"Slide {i+1}")
                is_placeholder = (suggested == f"Slide {i+1}")
                iss = self._new("slide_titles", i, "error",
                    f"Slide {i+1}: No title placeholder",
                    "Screen readers identify slides by title. This slide has none (WCAG 2.4.6).",
                    suggested_value=suggested)
                if self.should_auto("slide_titles") and not is_placeholder:
                    self._inject_title(slide, suggested)
                    iss.status = "auto_applied"
                    current = suggested
                self.issues.append(iss)

            else:
                current = raw_text(ts.text_frame._txBody)
                if not current:
                    body = self._body_text(slide)
                    suggested = (self._ai_title(i, body) or self._fallback_title(body) or f"Slide {i+1}")
                    is_placeholder = (suggested == f"Slide {i+1}")
                    iss = self._new("slide_titles", i, "error",
                        f"Slide {i+1}: Empty slide title",
                        "Title placeholder exists but contains no text.",
                        suggested_value=suggested)
                    if self.should_auto("slide_titles") and not is_placeholder:
                        set_title_text(ts, suggested)
                        iss.status = "auto_applied"
                        current = suggested
                    self.issues.append(iss)

            if current:
                if current in seen:
                    seen[current] += 1
                    deduped = f"{current} ({seen[current]})"
                    iss = self._new("duplicate_titles", i, "warning",
                        f"Slide {i+1}: Duplicate title",
                        f"Duplicate title '{current}' appears on multiple slides. Screen readers can't distinguish them.",
                        suggested_value=deduped)
                    if self.should_auto("duplicate_titles"):
                        ts2 = self._title_shape(slide)
                        if ts2: set_title_text(ts2, deduped)
                        iss.status = "auto_applied"
                    self.issues.append(iss)
                else:
                    seen[current] = 1

    def _check_broken_lists(self):
        for i, slide in enumerate(self.prs.slides):
            apply = self.should_auto("broken_lists")
            n = self._fix_broken_lists(slide, apply=apply)
            if n:
                iss = self._new("broken_lists", i, "warning",
                    f"Slide {i+1}: Broken list structure ({n} issue(s))",
                    "Gap paragraphs inside lists break list semantics for screen readers.",
                    suggested_value=str(n))
                iss.status = "auto_applied" if apply else "pending"
                self.issues.append(iss)

    def _check_image_alt_text(self):
        for i, slide in enumerate(self.prs.slides):
            context = self._title_text(slide) or self._body_text(slide, 200)
            for shape in slide.shapes:
                if shape.shape_type != MSO_SHAPE_TYPE.PICTURE: continue
                if get_alt_text(shape): continue

                blob, ct = self._image_data(shape)
                alt = self._ai_alt_text(blob, ct or "image/png", context) if (self.client and blob) else None
                is_placeholder = (alt is None)
                if not alt:
                    alt = f"Image on slide {i+1}" + (f" — {context[:40]}" if context else "")

                iss = self._new("image_alt_text", i, "error",
                    f'Slide {i+1}: Missing alt text on "{shape.name}"',
                    "Images require alt text for screen readers (WCAG 1.1.1).",
                    suggested_value=alt, shape_name=shape.name)
                if self.should_auto("image_alt_text") and not is_placeholder:
                    set_alt_text(shape, alt, alt)
                    iss.status = "auto_applied"
                self.issues.append(iss)

    def _check_tables(self):
        for i, slide in enumerate(self.prs.slides):
            slide_title = self._title_text(slide)
            for shape in slide.shapes:
                if shape.shape_type != MSO_SHAPE_TYPE.TABLE: continue

                n = unmerge_table_cells(shape)
                if n:
                    iss = self._new("table_unmerge", i, "error",
                        f'Slide {i+1}: {n} merged cell(s) in "{shape.name}"',
                        "Merged cells break table navigation for screen readers.",
                        shape_name=shape.name)
                    iss.status = "auto_applied"
                    self.issues.append(iss)

                if not get_alt_text(shape):
                    desc = describe_table(shape, slide_title)
                    iss = self._new("table_descriptions", i, "error",
                        f"Slide {i+1}: Table missing description",
                        "Tables need alt text to summarise their content (WCAG 1.1.1).",
                        suggested_value=desc, shape_name=shape.name)
                    if self.should_auto("table_descriptions"):
                        set_alt_text(shape, desc, desc)
                        iss.status = "auto_applied"
                    self.issues.append(iss)

                n = fix_empty_table_cells(shape.table)
                if n:
                    iss = self._new("table_empty_cells", i, "warning",
                        f'Slide {i+1}: {n} empty cell(s) in "{shape.name}"',
                        "Empty cells filled with an em-dash so screen readers announce them clearly.",
                        shape_name=shape.name)
                    iss.status = "auto_applied"
                    self.issues.append(iss)

    def _check_color_contrast(self):
        for i, slide in enumerate(self.prs.slides):
            for shape in slide.shapes:
                if not shape.has_text_frame: continue
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if not run.text.strip(): continue
                        color = get_run_color(run)
                        if color and is_too_light(color):
                            iss = self._new("color_contrast", i, "error",
                                f'Slide {i+1}: Low-contrast text in "{shape.name}"',
                                f"#{color} fails WCAG AA contrast (4.5:1 minimum).",
                                suggested_value=DARK_FALLBACK, shape_name=shape.name)
                            if self.should_auto("color_contrast"):
                                set_run_color(run, DARK_FALLBACK)
                                iss.status = "auto_applied"
                            self.issues.append(iss)

        # Always-human reminder for theme colors
        iss = self._new("theme_contrast", -1, "warning",
            "Theme-inherited color contrast — manual check required",
            "Colors inherited from the slide theme can't be read programmatically. "
            "Run Grackle Slides to verify remaining contrast issues.")
        iss.status = "pending"
        self.issues.append(iss)

    def _check_fine_print(self):
        for i, slide in enumerate(self.prs.slides):
            for shape in slide.shapes:
                if not shape.has_text_frame: continue
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if not run.text.strip(): continue
                        sz = run.font.size
                        if sz:
                            pt = sz / 12700
                            if 0 < pt < FINE_PRINT_PT:
                                iss = self._new("fine_print", i, "warning",
                                    f'Slide {i+1}: {pt:.0f}pt text in "{shape.name}"',
                                    f"Text is {pt:.0f}pt — increase to ≥{FINE_PRINT_PT}pt unless intentional.",
                                    shape_name=shape.name)
                                iss.status = "pending"
                                self.issues.append(iss)

    # ── FIX APPLICATION ────────────────────────────────────────────────────────

    def apply_fix(self, issue_id: str, action: str, value: str = None) -> dict:
        iss = next((i for i in self.issues if i.id == issue_id), None)
        if not iss:
            return {"error": "Issue not found"}

        if action == "skip":
            iss.status = "skipped"
            return {"status": "skipped", "slide_index": iss.slide_index}

        fix_value = value if (action == "edit" and value is not None) else iss.suggested_value
        success   = self._do_fix(iss, fix_value)
        iss.status = "accepted" if success else "pending"

        if success:
            self.save_pptx()

        return {
            "status":        iss.status,
            "slide_index":   iss.slide_index,
            "applied_value": fix_value,
        }

    def _do_fix(self, iss: Issue, value: str) -> bool:
        si = iss.slide_index

        if iss.check_type in ("slide_titles", "duplicate_titles"):
            if si < 0: return False
            slide = self.prs.slides[si]
            ts = self._title_shape(slide)
            if ts:
                set_title_text(ts, value); return True
            else:
                self._inject_title(slide, value); return True

        elif iss.check_type == "image_alt_text":
            if si < 0: return False
            slide = self.prs.slides[si]
            for shape in slide.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE and shape.name == iss.shape_name:
                    return bool(set_alt_text(shape, value, value))
            return False

        elif iss.check_type == "table_descriptions":
            if si < 0: return False
            slide = self.prs.slides[si]
            for shape in slide.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.TABLE and shape.name == iss.shape_name:
                    return bool(set_alt_text(shape, value, value))
            return False

        elif iss.check_type == "broken_lists":
            if si < 0: return False
            self._fix_broken_lists(self.prs.slides[si], apply=True)
            return True

        elif iss.check_type in ("fine_print", "theme_contrast", "color_contrast"):
            # Human acknowledges / marks as reviewed
            return True

        return False

    # ── THUMBNAILS ─────────────────────────────────────────────────────────────

    def generate_thumbnails(self):
        try:
            subprocess.run(
                ["libreoffice","--headless","--convert-to","pdf",
                 "--outdir", str(self.thumb_dir), str(self.pptx_path)],
                capture_output=True, timeout=90
            )
            pdf_path = self.thumb_dir / (self.pptx_path.stem + ".pdf")
            if not pdf_path.exists():
                pdfs = list(self.thumb_dir.glob("*.pdf"))
                if not pdfs: return []
                pdf_path = pdfs[0]

            from pdf2image import convert_from_path
            images = convert_from_path(str(pdf_path), dpi=120)
            self.thumbnail_paths = []
            for idx, img in enumerate(images):
                p = self.thumb_dir / f"slide_{idx}.png"
                img.save(str(p), "PNG")
                self.thumbnail_paths.append(str(p))
            return self.thumbnail_paths
        except Exception as e:
            print(f"[thumbnails] {e}")
            return []

    def regenerate_thumbnail(self, slide_index: int):
        if slide_index < 0: return
        try:
            subprocess.run(
                ["libreoffice","--headless","--convert-to","pdf",
                 "--outdir", str(self.thumb_dir), str(self.pptx_path)],
                capture_output=True, timeout=90
            )
            pdf_path = self.thumb_dir / (self.pptx_path.stem + ".pdf")
            if not pdf_path.exists():
                pdfs = list(self.thumb_dir.glob("*.pdf"))
                if not pdfs: return
                pdf_path = pdfs[0]

            from pdf2image import convert_from_path
            page = slide_index + 1
            images = convert_from_path(str(pdf_path), dpi=120, first_page=page, last_page=page)
            if images:
                p = self.thumb_dir / f"slide_{slide_index}.png"
                images[0].save(str(p), "PNG")
                if slide_index < len(self.thumbnail_paths):
                    self.thumbnail_paths[slide_index] = str(p)
        except Exception as e:
            print(f"[regen_thumb] {e}")

    def get_thumbnail_path(self, slide_index: int) -> Optional[str]:
        direct = self.thumb_dir / f"slide_{slide_index}.png"
        if direct.exists(): return str(direct)
        if slide_index < len(self.thumbnail_paths):
            p = self.thumbnail_paths[slide_index]
            if os.path.exists(p): return p
        return None

    # ── PERSISTENCE ────────────────────────────────────────────────────────────

    def save_pptx(self):
        self.prs.save(str(self.pptx_path))
        self.prs = Presentation(str(self.pptx_path))

    # ── SERIALISATION ──────────────────────────────────────────────────────────

    def to_dict(self) -> dict:
        slides = []
        for i in range(len(self.prs.slides)):
            slides.append({
                "index": i,
                "thumbnail_url": f"/api/thumbnail/{self.session_id}/{i}",
                "issue_count": sum(1 for iss in self.issues
                                   if iss.slide_index == i and iss.status == "pending"),
            })

        return {
            "session_id":   self.session_id,
            "filename":     self.original_path.name,
            "slide_count":  len(self.prs.slides),
            "slides":       slides,
            "issues":       [iss.to_dict() for iss in self.issues],
            "auto_applied": [iss.to_dict() for iss in self.issues if iss.status == "auto_applied"],
            "pending":      [iss.to_dict() for iss in self.issues if iss.status == "pending"],
            "checks":       {k: {**v, "enabled": self.should_auto(k)} for k, v in CHECKS.items()},
            "ai_enabled":   self.client is not None,
        }

    # ── PRIVATE HELPERS ────────────────────────────────────────────────────────

    def _new(self, check_type, slide_index, severity, title, description,
             suggested_value="", shape_name="") -> Issue:
        return Issue(
            id=str(uuid.uuid4()),
            slide_index=slide_index,
            check_type=check_type,
            severity=severity,
            title=title,
            description=description,
            suggested_value=suggested_value,
            shape_name=shape_name,
        )

    def _title_shape(self, slide):
        for shape in slide.shapes:
            if shape.has_text_frame and shape.is_placeholder:
                if shape.placeholder_format.type in (
                    PP_PLACEHOLDER.TITLE,
                    PP_PLACEHOLDER.CENTER_TITLE,
                    PP_PLACEHOLDER.VERTICAL_TITLE,
                ):
                    return shape
        return None

    def _title_text(self, slide) -> str:
        s = self._title_shape(slide)
        return raw_text(s.text_frame._txBody) if s else ""

    def _body_text(self, slide, max_chars=800) -> str:
        parts = []
        for shape in slide.shapes:
            if not shape.has_text_frame: continue
            if shape.is_placeholder and shape.placeholder_format.type in (
                PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE): continue
            t = raw_text(shape.text_frame._txBody)
            if t: parts.append(t)
        return " | ".join(parts)[:max_chars]

    def _image_data(self, shape):
        try: return shape.image.blob, shape.image.content_type
        except: return None, None

    def _ai_title(self, slide_index: int, body_text: str) -> Optional[str]:
        if not self.client or not body_text.strip(): return None
        try:
            msg = self.client.messages.create(
                model="claude-sonnet-4-6", max_tokens=30,
                messages=[{"role":"user","content":(
                    f"Slide {slide_index+1} of a presentation. "
                    "Write a short slide title (5 words max, Title Case). "
                    "Return ONLY the title.\n\n"
                    f"Content:\n{body_text[:600]}"
                )}]
            )
            return msg.content[0].text.strip().strip('"\'')
        except: return None

    def _ai_alt_text(self, blob: bytes, content_type: str, context: str = "") -> Optional[str]:
        if not self.client: return None
        mt = {"image/jpeg":"image/jpeg","image/jpg":"image/jpeg",
              "image/png":"image/png","image/gif":"image/gif",
              "image/webp":"image/webp"}.get((content_type or "").lower(),"image/png")
        b64 = base64.standard_b64encode(blob).decode()
        prompt = ("Write concise alt text (max 120 chars) for this presentation image. "
                  "Describe what it communicates. Do NOT start with 'Image of'. "
                  "Return only the alt text.")
        if context: prompt += f" Slide context: '{context}'."
        try:
            msg = self.client.messages.create(
                model="claude-sonnet-4-6", max_tokens=100,
                messages=[{"role":"user","content":[
                    {"type":"image","source":{"type":"base64","media_type":mt,"data":b64}},
                    {"type":"text","text":prompt},
                ]}]
            )
            return msg.content[0].text.strip()
        except: return None

    def _fallback_title(self, body_text: str) -> Optional[str]:
        if not body_text: return None
        for line in body_text.replace(" | ","\n").splitlines():
            line = line.strip().lstrip("•–—*-·").strip()
            if len(line) > 4:
                return (line[:50].rsplit(" ",1)[0]+"…") if len(line)>50 else line
        return None

    def _remove_trailing_paras(self, text_frame) -> int:
        removed = 0
        txBody = text_frame._txBody
        paras  = txBody.findall(atag("p"))
        while len(paras) > 1:
            last = paras[-1]
            if raw_text(last): break
            txBody.remove(last)
            paras = txBody.findall(atag("p"))
            removed += 1
        return removed

    def _fix_broken_lists(self, slide, apply=True) -> int:
        count = 0
        for shape in slide.shapes:
            if not shape.has_text_frame: continue
            paras = shape.text_frame.paragraphs
            for j, para in enumerate(paras):
                pPr    = para._p.find(qn("a:pPr"))
                if pPr is None: continue
                buNone = pPr.find(qn("a:buNone"))
                if buNone is not None and 0 < j < len(paras)-1:
                    prev_bullet = self._has_bullet(paras[j-1]._p)
                    next_bullet = self._has_bullet(paras[j+1]._p)
                    if prev_bullet and next_bullet and not raw_text(para._p):
                        if apply: pPr.remove(buNone)
                        count += 1
        return count

    def _has_bullet(self, p_el) -> bool:
        pPr = p_el.find(qn("a:pPr"))
        return pPr is None or pPr.find(qn("a:buNone")) is None

    def _inject_title(self, slide, text: str):
        spTree  = slide.shapes._spTree
        ids     = [int(el.get("id")) for el in spTree.iter()
                   if el.get("id") and el.get("id").isdigit()]
        new_id  = max(ids, default=100) + 1
        safe    = xml_escape(text)
        sp_xml  = (
            f'<p:sp xmlns:p="{PML}" xmlns:a="{DML}">'
            f'<p:nvSpPr>'
            f'  <p:cNvPr id="{new_id}" name="Title {new_id}" descr="{safe}" title="{safe}"/>'
            f'  <p:cNvSpPr><a:spLocks noGrp="1"/></p:cNvSpPr>'
            f'  <p:nvPr><p:ph type="title"/></p:nvPr>'
            f'</p:nvSpPr>'
            f'<p:spPr><a:xfrm>'
            f'  <a:off x="457200" y="274638"/><a:ext cx="8229600" cy="1143000"/>'
            f'</a:xfrm><a:prstGeom prst="rect"><a:avLst/></a:prstGeom></p:spPr>'
            f'<p:txBody><a:bodyPr/><a:lstStyle/>'
            f'<a:p><a:r><a:rPr lang="{DEFAULT_LANGUAGE}" dirty="0" b="1"/>'
            f'<a:t>{safe}</a:t></a:r></a:p>'
            f'</p:txBody></p:sp>'
        )
        spTree.append(etree.fromstring(sp_xml))
